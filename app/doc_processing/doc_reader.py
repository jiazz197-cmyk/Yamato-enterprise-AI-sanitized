import json
import logging
import os
import subprocess
import sys
import tempfile
import time
import uuid
from datetime import datetime
from io import BytesIO
from pathlib import Path
from typing import Dict, List, Tuple, Union

import html_text
import pandas as pd
import pdfplumber
import readability
from docx import Document as DocxDocument
from pptx import Presentation
from langchain_core.documents import Document

import langchain_compat  # noqa: F401

from .exceptions import DocumentParseError
from .text_splitter import TagGenerator, TokenAwareTextSplitter

logger = logging.getLogger(__name__)

try:
    from paddleocr import PaddleOCR
except ImportError:  # pragma: no cover
    PaddleOCR = None  # type: ignore


def file_to_stream(file_path: Union[str, os.PathLike], keep_filename: bool = True) -> BytesIO:
    """将文件转换成 BytesIO，供解析器消费"""
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"文件不存在: {file_path}")

    with open(file_path, "rb") as f:
        file_content = f.read()

    stream = BytesIO(file_content)
    if keep_filename:
        stream.name = os.path.basename(file_path)
    return stream


def file_to_bytesio(file_path: Union[str, os.PathLike]) -> BytesIO:
    """兼容旧接口"""
    return file_to_stream(file_path, keep_filename=True)


class LibreOfficeConverter:
    """使用 LibreOffice 将 .doc 转换为 .docx"""

    def __init__(self, soffice_path: str | None = None):
        if soffice_path is not None:
            self.soffice_path = soffice_path
        elif os.name == "nt":
            self.soffice_path = r"C:\Program Files\LibreOffice\program\soffice.exe"
        elif sys.platform == "darwin":
            default_path = "/Applications/LibreOffice.app/Contents/MacOS/soffice"
            self.soffice_path = default_path if os.path.exists(default_path) else "soffice"
        else:
            self.soffice_path = "soffice"
        self._ensure_available()

    def _ensure_available(self):
        try:
            subprocess.run(
                [self.soffice_path, "--version"],
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE,
                check=True,
                timeout=5,
            )
        except (subprocess.SubprocessError, FileNotFoundError):
            if sys.platform == "darwin" and self.soffice_path == "soffice":
                fallback = "/Applications/LibreOffice.app/Contents/MacOS/soffice"
                if os.path.exists(fallback):
                    self.soffice_path = fallback
                    subprocess.run([self.soffice_path, "--version"], check=True, timeout=5)
                    return
            raise RuntimeError(f"未检测到LibreOffice，请确认路径：{self.soffice_path}")

    def convert_to_docx(self, file_input: Union[str, os.PathLike, bytes, BytesIO]) -> BytesIO:
        if isinstance(file_input, (str, os.PathLike)) and not str(file_input).lower().endswith(".doc"):
            raise ValueError("仅支持 .doc 转换")

        with tempfile.NamedTemporaryFile(suffix=".doc", delete=False) as tmp_doc:
            if isinstance(file_input, (str, os.PathLike)):
                with open(file_input, "rb") as src:
                    tmp_doc.write(src.read())
            elif isinstance(file_input, BytesIO):
                file_input.seek(0)
                tmp_doc.write(file_input.read())
            else:
                tmp_doc.write(file_input if isinstance(file_input, bytes) else bytes(file_input))
            tmp_doc.flush()
            os.fsync(tmp_doc.fileno())
            doc_path = tmp_doc.name

        with tempfile.TemporaryDirectory() as out_dir:
            try:
                proc = subprocess.run(
                    [
                        self.soffice_path,
                        "--headless",
                        "--convert-to",
                        "docx:Office Open XML Text",
                        "--outdir",
                        out_dir,
                        doc_path,
                    ],
                    stdout=subprocess.PIPE,
                    stderr=subprocess.PIPE,
                    check=True,
                    text=True,
                )
                if proc.returncode != 0:
                    raise RuntimeError(proc.stderr)
                result_path = os.path.join(out_dir, Path(doc_path).with_suffix(".docx").name)
                if not os.path.exists(result_path):
                    raise RuntimeError("转换后文件不存在")
                with open(result_path, "rb") as f:
                    return BytesIO(f.read())
            finally:
                if os.path.exists(doc_path):
                    for _ in range(3):
                        try:
                            os.remove(doc_path)
                            break
                        except PermissionError:
                            time.sleep(0.1)
                    else:
                        logger.warning("无法删除临时文件：%s", doc_path)


class PdfParser:
    """PDF 解析，必要时使用 OCR"""

    def __init__(self, save_dir: str = "paddle_ocr_images", enable_ocr: bool = True):
        self.save_dir = save_dir
        os.makedirs(save_dir, exist_ok=True)
        self.enable_ocr = enable_ocr and PaddleOCR is not None
        self.ocr = PaddleOCR(use_angle_cls=True, lang="ch") if self.enable_ocr else None

    def __call__(self, file_input: Union[str, bytes, os.PathLike, BytesIO]) -> Tuple[str, List[Dict]]:
        temp_path = None
        try:
            if isinstance(file_input, (str, os.PathLike)):
                file_path = str(file_input)
            else:
                temp_path = os.path.join(self.save_dir, f"temp_{uuid.uuid4().hex}.pdf")
                with open(temp_path, "wb") as temp_file:
                    if isinstance(file_input, bytes):
                        temp_file.write(file_input)
                    else:
                        file_input.seek(0)
                        temp_file.write(file_input.read())
                file_path = temp_path
            text = self._extract_text(file_path)
            return text, []
        finally:
            if temp_path and os.path.exists(temp_path):
                try:
                    os.remove(temp_path)
                except OSError:
                    pass

    def _extract_text(self, file_path: str) -> str:
        chunks: List[str] = []
        with pdfplumber.open(file_path) as pdf:
            for idx, page in enumerate(pdf.pages):
                page_text = page.extract_text()
                if page_text and page_text.strip():
                    chunks.append(page_text)
                    continue
                if self.ocr is None:
                    continue
                try:
                    image_path = os.path.join(self.save_dir, f"page_{idx + 1}.png")
                    page.to_image(resolution=300).original.save(image_path)
                    result = self.ocr.ocr(image_path, cls=True) or []
                    ocr_text = "\n".join(seg[1][0] for line in result for seg in line)
                    chunks.append(ocr_text)
                finally:
                    if os.path.exists(image_path):
                        os.remove(image_path)
        return "\n".join(chunks)


class DocxParser:
    """DOCX 解析器"""

    def __call__(self, file_input: Union[str, bytes, os.PathLike, BytesIO]) -> Tuple[str, List[Dict]]:
        doc = self._load_doc(file_input)
        paragraphs = [para.text for para in doc.paragraphs]
        tables = [self._table_to_dict(table) for table in doc.tables]
        return "\n".join(paragraphs), tables

    def _load_doc(self, file_input):
        if isinstance(file_input, (str, os.PathLike)):
            return DocxDocument(file_input)
        if isinstance(file_input, bytes):
            return DocxDocument(BytesIO(file_input))
        if isinstance(file_input, BytesIO):
            file_input.seek(0)
            return DocxDocument(file_input)
        raise ValueError("不支持的输入类型")

    def _table_to_dict(self, table) -> Dict:
        headers = [cell.text.strip() for cell in table.rows[0].cells] if table.rows else []
        rows = []
        for row in table.rows[1:]:
            rows.append([cell.text.strip() for cell in row.cells])
        return {"headers": headers, "rows": rows}

    def _table_to_text(self, table_dict: Dict) -> str:
        headers = table_dict.get("headers", [])
        rows = table_dict.get("rows", [])
        output = ["\t".join(headers)]
        for row in rows:
            output.append("\t".join(row))
        return "\n".join(output)


class DocParser:
    """DOC 文件解析（转 DOCX）"""

    def __init__(self):
        self.converter = LibreOfficeConverter()
        self.docx_parser = DocxParser()

    def __call__(self, file_input: Union[str, bytes, os.PathLike, BytesIO]) -> Tuple[str, List[Dict]]:
        docx_bytes = self.converter.convert_to_docx(file_input)
        return self.docx_parser(docx_bytes)


class ExcelParser:
    """Excel 解析为纯文本"""

    def __call__(
        self,
        file_input: Union[str, bytes, os.PathLike, BytesIO],
        sheet_idx: int = 0,
    ) -> Tuple[str, List[Dict]]:
        temp_path = None
        try:
            if isinstance(file_input, (str, os.PathLike)):
                path = str(file_input)
            else:
                suffix = ".xlsx"
                temp_path = f"temp_excel_{uuid.uuid4().hex}{suffix}"
                with open(temp_path, "wb") as tmp:
                    if isinstance(file_input, bytes):
                        tmp.write(file_input)
                    else:
                        file_input.seek(0)
                        tmp.write(file_input.read())
                path = temp_path

            df = pd.read_excel(path, sheet_name=sheet_idx, header=None)
            df = df.fillna("")
            text = "\n".join("\t".join(map(str, row)) for row in df.values.tolist())
            table = {"headers": df.iloc[0].tolist() if not df.empty else [], "rows": df.iloc[1:].values.tolist()}
            return text, [table]
        finally:
            if temp_path and os.path.exists(temp_path):
                os.remove(temp_path)


class PptParser:
    """PPT/PPTX 解析"""

    def __call__(self, file_input: Union[str, bytes, os.PathLike, BytesIO]) -> Tuple[str, List[Dict]]:
        presentation = self._load_presentation(file_input)
        texts: List[str] = []
        for slide_idx, slide in enumerate(presentation.slides, start=1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text.append(shape.text)
            if slide_text:
                texts.append(f"Slide {slide_idx}:\n" + "\n".join(slide_text))
        return "\n\n".join(texts), []

    def _load_presentation(self, file_input):
        if isinstance(file_input, (str, os.PathLike)):
            return Presentation(file_input)
        data = file_input if isinstance(file_input, bytes) else file_input.getvalue()
        return Presentation(BytesIO(data))


class JsonParser:
    def __call__(self, file_input: Union[str, bytes, os.PathLike, BytesIO]) -> Tuple[str, List[Dict]]:
        data = self._load_json(file_input)
        return json.dumps(data, ensure_ascii=False, indent=2), []

    def _load_json(self, file_input):
        if isinstance(file_input, (str, os.PathLike)):
            with open(file_input, "r", encoding="utf-8") as f:
                return json.load(f)
        if isinstance(file_input, bytes):
            return json.loads(file_input.decode("utf-8"))
        if isinstance(file_input, BytesIO):
            file_input.seek(0)
            return json.load(file_input)
        raise ValueError("不支持的输入类型")


class TxtParser:
    def __call__(self, file_input: Union[str, bytes, os.PathLike, BytesIO]) -> Tuple[str, List[Dict]]:
        if isinstance(file_input, (str, os.PathLike)):
            with open(file_input, "r", encoding="utf-8", errors="ignore") as f:
                return f.read(), []
        if isinstance(file_input, bytes):
            return file_input.decode("utf-8", errors="ignore"), []
        if isinstance(file_input, BytesIO):
            file_input.seek(0)
            return file_input.read().decode("utf-8", errors="ignore"), []
        raise ValueError("不支持的输入类型")


class HtmlParser:
    def __call__(self, file_input: Union[str, bytes, os.PathLike, BytesIO]) -> Tuple[str, List[Dict]]:
        html = self._load_html(file_input)
        doc = readability.Document(html)
        text = html_text.extract_text(doc.summary(html_partial=True))
        metadata = {
            "title": doc.short_title(),
            "processed_at": datetime.now().isoformat(),
        }
        return f"{metadata['title']}\n{text}", []

    def _load_html(self, file_input) -> str:
        if isinstance(file_input, (str, os.PathLike)):
            with open(file_input, "r", encoding="utf-8", errors="ignore") as f:
                return f.read()
        if isinstance(file_input, bytes):
            return file_input.decode("utf-8", errors="ignore")
        if isinstance(file_input, BytesIO):
            file_input.seek(0)
            return file_input.read().decode("utf-8", errors="ignore")
        raise ValueError("不支持的输入类型")


class DocumentProcessor:
    """统一文档解析器"""

    def __init__(self):
        self.parsers = {
            "pdf": PdfParser(),
            "docx": DocxParser(),
            "doc": DocParser(),
            "xlsx": ExcelParser(),
            "xls": ExcelParser(),
            "pptx": PptParser(),
            "ppt": PptParser(),
            "html": HtmlParser(),
            "htm": HtmlParser(),
            "json": JsonParser(),
            "txt": TxtParser(),
            "md": TxtParser(),
        }
        self.processed_files = set()
        self.failed_files: Dict[str, str] = {}

    def get_file_extension(self, file_input) -> str:
        if isinstance(file_input, BytesIO) and hasattr(file_input, "name"):
            file_name = file_input.name
        elif isinstance(file_input, (str, os.PathLike)):
            file_name = os.path.basename(str(file_input))
        elif hasattr(file_input, "name"):
            file_name = file_input.name
        else:
            logger.debug(f"无法获取文件扩展名: 类型={type(file_input)}")
            return "unknown"

        # 处理空字符串或无效文件名
        if not file_name or not file_name.strip():
            logger.warning(f"文件名为空: {type(file_input)}")
            return "unknown"

        if "." in file_name:
            ext = os.path.splitext(file_name)[1][1:].lower()
            # 处理文件名以点号结尾的情况（如 "test."）或空扩展名
            if not ext:
                logger.warning(f"文件名以点号结尾或无扩展名: {file_name}")
                return "unknown"
            return ext
        
        logger.debug(f"文件名无扩展名: {file_name}")
        return "unknown"

    def extract_metadata(self, file_input: Union[str, bytes, os.PathLike], text: str) -> Dict:
        metadata = {
            "source": "",
            "file_type": "",
            "file_size": 0,
            "created_time": datetime.now().isoformat(),
            "modified_time": datetime.now().isoformat(),
            "title": "",
            "author": "未知",
            "processing_time": datetime.now().isoformat(),
            "chunk_count": 0,
        }

        try:
            if isinstance(file_input, (str, os.PathLike)):
                file_stats = os.stat(file_input)
                metadata.update(
                    {
                        "source": os.path.basename(file_input),
                        "file_type": os.path.splitext(file_input)[1][1:].lower(),
                        "file_size": file_stats.st_size,
                        "created_time": datetime.fromtimestamp(file_stats.st_ctime).isoformat(),
                        "modified_time": datetime.fromtimestamp(file_stats.st_mtime).isoformat(),
                        "title": os.path.splitext(os.path.basename(file_input))[0],
                    }
                )
            else:
                file_name = getattr(file_input, "name", "uploaded_file")
                metadata.update(
                    {
                        "source": file_name,
                        "file_type": file_name.split(".")[-1].lower() if "." in file_name else "unknown",
                        "title": file_name.split(".")[0] if "." in file_name else file_name,
                    }
                )

            self._extract_content_metadata(text, metadata)
        except Exception as exc:
            logger.warning("提取元数据失败: %s", exc, exc_info=True)

        return metadata

    def _extract_content_metadata(self, text: str, metadata: Dict):
        import re

        first_page = text[:2000]
        author_patterns = [r"作者[:：]\s*(\S+)", r"by\s*(\S+)", r"author[:：]\s*(\S+)"]
        for pattern in author_patterns:
            if match := re.search(pattern, first_page, re.IGNORECASE):
                metadata["author"] = match.group(1).strip()
                break

        date_patterns = [r"日期[:：]\s*(\d{4}-\d{2}-\d{2})", r"date[:：]\s*(\d{4}-\d{2}-\d{2})", r"\b(\d{4}年\d{1,2}月\d{1,2}日)\b"]
        for pattern in date_patterns:
            if match := re.search(pattern, first_page, re.IGNORECASE):
                metadata["date"] = match.group(1).strip()
                break

    def process_document(
        self,
        file_input,
        text_splitter: TokenAwareTextSplitter,
        tag_generator: TagGenerator = None,
        num_tags: int = 5,
    ) -> List[Document]:
        try:
            file_ext = self.get_file_extension(file_input)
            
            # 获取文件名用于调试
            if isinstance(file_input, BytesIO):
                file_name = getattr(file_input, "name", "<BytesIO without name>")
            elif isinstance(file_input, (str, os.PathLike)):
                file_name = os.path.basename(str(file_input))
            else:
                file_name = str(file_input)
            
            if file_ext not in self.parsers:
                supported = ", ".join(sorted(self.parsers.keys()))
                raise DocumentParseError(
                    f"不支持的文件类型: '{file_ext}' (文件: {file_name}, 支持的格式: {supported})"
                )

            parser = self.parsers[file_ext]
            text, tables = parser(file_input)

            if tables:
                table_texts = []
                for table in tables:
                    if hasattr(parser, "_table_to_text"):
                        table_texts.append(parser._table_to_text(table))
                    else:
                        table_texts.append(str(table))
                full_text = f"{text}\n\n{chr(10).join(table_texts)}"
            else:
                full_text = text

            metadata = self.extract_metadata(file_input, full_text)
            doc = Document(page_content=full_text, metadata=metadata)
            split_texts = text_splitter.split_text(doc.page_content)

            chunks: List[Document] = []
            for text_chunk in split_texts:
                chunk_metadata = {
                    **doc.metadata,
                    "token_count": text_splitter.count_tokens(text_chunk),
                }
                if tag_generator and text_chunk.strip():
                    chunk_metadata["tags"] = tag_generator.extract_tags(text_chunk, num_tags=num_tags)
                chunks.append(Document(page_content=text_chunk, metadata=chunk_metadata))

            metadata["chunk_count"] = len(chunks)
            return chunks
        except DocumentParseError:
            raise
        except Exception as exc:
            file_name = str(file_input) if isinstance(file_input, (str, os.PathLike)) else getattr(file_input, "name", "file")
            self.failed_files[file_name] = str(exc)
            logger.exception("处理文档失败: %s", file_name)
            raise DocumentParseError(str(exc)) from exc

    def convert_to_excel_stream(self, text: str, tables: List[Dict]) -> BytesIO:
        """保留老逻辑，便于外部导出"""
        try:
            output = BytesIO()
            with pd.ExcelWriter(output, engine="openpyxl") as writer:
                text_df = pd.DataFrame({"文本内容": text.split("\n")})
                text_df.to_excel(writer, sheet_name="文本内容", index=False)

                for i, table in enumerate(tables):
                    if isinstance(table, dict) and "headers" in table and "rows" in table:
                        headers = table["headers"]
                        rows = [row["cells"] for row in table.get("rows", [])]
                    else:
                        headers = [f"列{j + 1}" for j in range(len(table[0]))]
                        rows = table

                    sheet_name = f"表格_{i + 1}"[:31]
                    table_df = pd.DataFrame(rows, columns=headers)
                    table_df.to_excel(writer, sheet_name=sheet_name, index=False)

            output.seek(0)
            return output
        except Exception as exc:
            logger.exception("导出Excel失败")
            raise DocumentParseError(f"导出Excel失败: {exc}") from exc

